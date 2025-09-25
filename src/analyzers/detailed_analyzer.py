import os
import io
import base64
from typing import Dict, List, Any, Tuple
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageStat
from colorthief import ColorThief
import webcolors
import numpy as np
import cv2


class DetailedPPTAnalyzer:
    """Enhanced PPT analyzer that extracts comprehensive details including colors, images, and formatting."""
    
    def __init__(self):
        self.extracted_images = []
        self.color_palette = []
        self.fonts_used = set()
        
    def parse_ppt_detailed(self, file_path: str) -> Dict[str, Any]:
        """Parse PPT with comprehensive details extraction."""
        prs = Presentation(file_path)
        
        analysis_result = {
            "presentation_info": self._get_presentation_info(prs),
            "slides": [],
            "global_analysis": {
                "total_slides": len(prs.slides),
                "color_palette": [],
                "fonts_used": [],
                "image_count": 0,
                "slide_layouts": []
            }
        }
        
        # Process each slide
        for i, slide in enumerate(prs.slides, start=1):
            slide_data = self._analyze_slide(slide, i)
            analysis_result["slides"].append(slide_data)
            
            # Update global statistics
            analysis_result["global_analysis"]["image_count"] += slide_data["image_count"]
            analysis_result["global_analysis"]["slide_layouts"].append(slide_data["layout"])
        
        # Extract global color palette and fonts
        analysis_result["global_analysis"]["color_palette"] = list(self.color_palette)
        analysis_result["global_analysis"]["fonts_used"] = list(self.fonts_used)
        
        return analysis_result
    
    def _get_presentation_info(self, prs: Presentation) -> Dict[str, Any]:
        """Extract presentation-level metadata."""
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        return {
            "slide_dimensions": {
                "width": slide_width,
                "height": slide_height,
                "width_inches": slide_width / 914400,  # Convert EMU to inches
                "height_inches": slide_height / 914400,
                "aspect_ratio": round(slide_width / slide_height, 2)
            },
            "slide_master_count": len(prs.slide_masters),
            "slide_layouts_count": sum(len(master.slide_layouts) for master in prs.slide_masters)
        }
    
    def _analyze_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """Analyze individual slide with comprehensive details."""
        slide_data = {
            "slide_number": slide_number,
            "layout": slide.slide_layout.name,
            "background": self._analyze_background(slide),
            "shapes": [],
            "texts": [],
            "images": [],
            "colors": [],
            "fonts": [],
            "image_count": 0,
            "text_count": 0,
            "shape_count": len(slide.shapes)
        }
        
        # Analyze each shape
        for shape_idx, shape in enumerate(slide.shapes):
            shape_analysis = self._analyze_shape(shape, shape_idx)
            slide_data["shapes"].append(shape_analysis)
            
            # Collect texts
            if shape_analysis["text"]:
                slide_data["texts"].extend(shape_analysis["text"])
                slide_data["text_count"] += len(shape_analysis["text"])
            
            # Collect images
            if shape_analysis["type"] == "picture":
                slide_data["images"].append(shape_analysis)
                slide_data["image_count"] += 1
            
            # Collect colors and fonts
            slide_data["colors"].extend(shape_analysis["colors"])
            slide_data["fonts"].extend(shape_analysis["fonts"])
        
        # Remove duplicates and update global collections
        slide_data["colors"] = list(set(slide_data["colors"]))
        slide_data["fonts"] = list(set(slide_data["fonts"]))
        
        self.color_palette.extend(slide_data["colors"])
        self.fonts_used.update(slide_data["fonts"])
        
        return slide_data
    
    def _analyze_background(self, slide) -> Dict[str, Any]:
        """Analyze slide background."""
        background_info = {
            "type": "unknown",
            "color": None,
            "has_image": False,
            "gradient": False
        }
        
        try:
            background = slide.background
            if hasattr(background, 'fill'):
                fill = background.fill
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    background_info["color"] = self._extract_color(fill.fore_color)
                    background_info["type"] = "solid_color"
        except:
            pass
        
        return background_info
    
    def _analyze_shape(self, shape, shape_idx: int) -> Dict[str, Any]:
        """Analyze individual shape with all details."""
        shape_data = {
            "index": shape_idx,
            "type": self._get_shape_type(shape),
            "name": getattr(shape, 'name', f'Shape_{shape_idx}'),
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            },
            "text": [],
            "colors": [],
            "fonts": [],
            "image_info": None
        }
        
        # Analyze text content
        try:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                text_analysis = self._analyze_text(shape)
                shape_data["text"] = text_analysis["texts"]
                shape_data["fonts"].extend(text_analysis["fonts"])
                shape_data["colors"].extend(text_analysis["colors"])
        except Exception as e:
            # If text analysis fails, just skip it
            pass
        
        # Analyze fill colors
        if hasattr(shape, 'fill'):
            fill_color = self._extract_fill_color(shape.fill)
            if fill_color:
                shape_data["colors"].append(fill_color)
        
        # Analyze line colors
        if hasattr(shape, 'line'):
            line_color = self._extract_line_color(shape.line)
            if line_color:
                shape_data["colors"].append(line_color)
        
        # Analyze images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_data["image_info"] = self._analyze_image(shape)
            if shape_data["image_info"]:
                shape_data["colors"].extend(shape_data["image_info"]["dominant_colors"])
        
        return shape_data
    
    def _analyze_text(self, shape) -> Dict[str, Any]:
        """Analyze text formatting and content."""
        text_data = {
            "texts": [],
            "fonts": [],
            "colors": []
        }
        
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        text_info = {
                            "content": run.text.strip(),
                            "font": {
                                "name": run.font.name if run.font.name else "Unknown",
                                "size": run.font.size.pt if run.font.size else None,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": run.font.underline
                            }
                        }
                        
                        # Extract text color
                        try:
                            if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                color = self._rgb_to_hex(run.font.color.rgb)
                                if color:
                                    text_info["color"] = color
                                    text_data["colors"].append(color)
                        except:
                            pass
                        
                        text_data["texts"].append(text_info)
                        text_data["fonts"].append(text_info["font"]["name"])
        
        return text_data
    
    def _analyze_image(self, shape) -> Dict[str, Any]:
        """Analyze image properties and extract dominant colors."""
        try:
            image = shape.image
            image_stream = io.BytesIO(image.blob)
            
            # Open with PIL
            pil_image = Image.open(image_stream)
            
            # Basic image info
            image_info = {
                "format": pil_image.format,
                "mode": pil_image.mode,
                "size": pil_image.size,
                "filename": getattr(image, 'filename', 'unknown'),
                "dominant_colors": [],
                "average_color": None,
                "brightness": None,
                "base64": None
            }
            
            # Convert to RGB if necessary
            if pil_image.mode != 'RGB':
                rgb_image = pil_image.convert('RGB')
            else:
                rgb_image = pil_image
            
            # Calculate average color and brightness
            stat = ImageStat.Stat(rgb_image)
            avg_color = tuple(int(c) for c in stat.mean)
            image_info["average_color"] = self._rgb_to_hex_tuple(avg_color)
            image_info["brightness"] = sum(avg_color) / 3
            
            # Extract dominant colors using ColorThief
            try:
                # Save temporary image for ColorThief
                temp_stream = io.BytesIO()
                rgb_image.save(temp_stream, format='PNG')
                temp_stream.seek(0)
                
                color_thief = ColorThief(temp_stream)
                dominant_colors = color_thief.get_palette(color_count=5, quality=10)
                
                image_info["dominant_colors"] = [
                    self._rgb_to_hex_tuple(color) for color in dominant_colors
                ]
            except:
                # Fallback method using OpenCV
                image_info["dominant_colors"] = self._extract_colors_opencv(rgb_image)
            
            # Convert to base64 for web display
            buffer = io.BytesIO()
            rgb_image.save(buffer, format='PNG')
            image_info["base64"] = base64.b64encode(buffer.getvalue()).decode()
            
            return image_info
            
        except Exception as e:
            return {"error": str(e)}
    
    def _extract_colors_opencv(self, pil_image) -> List[str]:
        """Extract dominant colors using OpenCV as fallback."""
        try:
            # Convert PIL to OpenCV format
            opencv_image = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
            
            # Reshape image to be a list of pixels
            data = opencv_image.reshape((-1, 3))
            data = np.float32(data)
            
            # Apply k-means clustering
            criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 20, 1.0)
            k = 5
            _, labels, centers = cv2.kmeans(data, k, None, criteria, 10, cv2.KMEANS_RANDOM_CENTERS)
            
            # Convert centers back to RGB and then to hex
            centers = np.uint8(centers)
            dominant_colors = []
            for center in centers:
                # Convert BGR back to RGB
                rgb_color = (int(center[2]), int(center[1]), int(center[0]))
                dominant_colors.append(self._rgb_to_hex_tuple(rgb_color))
            
            return dominant_colors
        except:
            return []
    
    def _get_shape_type(self, shape) -> str:
        """Get human-readable shape type."""
        try:
            shape_types = {
                MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
                MSO_SHAPE_TYPE.CALLOUT: "callout",
                MSO_SHAPE_TYPE.CHART: "chart",
                MSO_SHAPE_TYPE.FREEFORM: "freeform",
                MSO_SHAPE_TYPE.GROUP: "group",
                MSO_SHAPE_TYPE.LINE: "line",
                MSO_SHAPE_TYPE.MEDIA: "media",
                MSO_SHAPE_TYPE.PICTURE: "picture",
                MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
                MSO_SHAPE_TYPE.TABLE: "table",
                MSO_SHAPE_TYPE.TEXT_BOX: "text_box"
            }
            
            # Add additional types if they exist
            if hasattr(MSO_SHAPE_TYPE, 'COMMENT'):
                shape_types[MSO_SHAPE_TYPE.COMMENT] = "comment"
            if hasattr(MSO_SHAPE_TYPE, 'LINKED_OLE_OBJECT'):
                shape_types[MSO_SHAPE_TYPE.LINKED_OLE_OBJECT] = "linked_ole_object"
            if hasattr(MSO_SHAPE_TYPE, 'LINKED_PICTURE'):
                shape_types[MSO_SHAPE_TYPE.LINKED_PICTURE] = "linked_picture"
            if hasattr(MSO_SHAPE_TYPE, 'OLE_OBJECT'):
                shape_types[MSO_SHAPE_TYPE.OLE_OBJECT] = "ole_object"
            if hasattr(MSO_SHAPE_TYPE, 'SCRIPT_ANCHOR'):
                shape_types[MSO_SHAPE_TYPE.SCRIPT_ANCHOR] = "script_anchor"
                
            return shape_types.get(shape.shape_type, f"unknown_type_{shape.shape_type}")
        except Exception as e:
            return f"error_getting_type_{str(e)[:20]}"
    
    def _extract_color(self, color_obj) -> str:
        """Extract color from color object."""
        try:
            if color_obj is None:
                return None
            if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
                return self._rgb_to_hex(color_obj.rgb)
        except:
            pass
        return None
    
    def _extract_fill_color(self, fill) -> str:
        """Extract fill color from shape."""
        try:
            if hasattr(fill, 'fore_color') and fill.fore_color:
                return self._extract_color(fill.fore_color)
        except:
            pass
        return None
    
    def _extract_line_color(self, line) -> str:
        """Extract line color from shape."""
        try:
            if hasattr(line, 'color') and line.color:
                return self._extract_color(line.color)
        except:
            pass
        return None
    
    def _rgb_to_hex(self, rgb_color) -> str:
        """Convert RGB color to hex string."""
        try:
            if rgb_color is None:
                return None
            if isinstance(rgb_color, RGBColor):
                return f"#{rgb_color.r:02x}{rgb_color.g:02x}{rgb_color.b:02x}"
            elif hasattr(rgb_color, 'r') and hasattr(rgb_color, 'g') and hasattr(rgb_color, 'b'):
                return f"#{rgb_color.r:02x}{rgb_color.g:02x}{rgb_color.b:02x}"
            return str(rgb_color)
        except:
            return None
    
    def _rgb_to_hex_tuple(self, rgb_tuple: Tuple[int, int, int]) -> str:
        """Convert RGB tuple to hex string."""
        return f"#{rgb_tuple[0]:02x}{rgb_tuple[1]:02x}{rgb_tuple[2]:02x}"
    
    def get_color_name(self, hex_color: str) -> str:
        """Get closest color name for a hex color."""
        try:
            rgb = webcolors.hex_to_rgb(hex_color)
            return webcolors.rgb_to_name(rgb)
        except ValueError:
            # Find closest color
            try:
                closest_name = self._closest_color(hex_color)
                return closest_name
            except:
                return "Unknown"
    
    def _closest_color(self, hex_color: str) -> str:
        """Find the closest named color."""
        try:
            rgb_color = webcolors.hex_to_rgb(hex_color)
            min_colors = {}
            
            for key, name in webcolors.CSS3_HEX_TO_NAMES.items():
                r_c, g_c, b_c = webcolors.hex_to_rgb(key)
                rd = (r_c - rgb_color[0]) ** 2
                gd = (g_c - rgb_color[1]) ** 2
                bd = (b_c - rgb_color[2]) ** 2
                min_colors[(rd + gd + bd)] = name
            
            return min_colors[min(min_colors.keys())]
        except:
            return "Unknown"

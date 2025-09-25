import os
import sys
import tempfile
import traceback
from flask import Flask, request, jsonify, render_template, redirect, url_for, flash
from werkzeug.utils import secure_filename

# Add parent directory to path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

app = Flask(__name__, 
           static_folder='../static',
           template_folder='../templates')

app.secret_key = os.getenv('FLASK_SECRET_KEY', 'minimal-ppt-analyzer')

# Extremely conservative limits for Vercel
ALLOWED_EXTENSIONS = {'pptx', 'ppt'}
MAX_FILE_SIZE = 500 * 1024  # 500KB - very conservative for Vercel
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

def allowed_file(filename):
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def validate_request_size():
    """Validate request size before processing."""
    content_length = request.content_length
    if content_length and content_length > MAX_FILE_SIZE:
        return False, f'Request too large: {content_length} bytes. Maximum: {MAX_FILE_SIZE} bytes (500KB)'
    return True, None

def minimal_ppt_analysis(file_path):
    """Minimal PPT analysis that won't crash serverless."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        
        # Very basic analysis to avoid memory issues
        result = {
            'filename': os.path.basename(file_path),
            'total_slides': len(prs.slides),
            'slide_dimensions': {
                'width': prs.slide_width,
                'height': prs.slide_height
            },
            'slides': []
        }
        
        # Process only first 5 slides to avoid memory issues
        for i, slide in enumerate(prs.slides[:5], 1):
            slide_info = {
                'slide_number': i,
                'shapes_count': len(slide.shapes),
                'layout_name': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
            }
            result['slides'].append(slide_info)
        
        return result
        
    except Exception as e:
        return {'error': f'Analysis failed: {str(e)}'}

@app.route('/')
def index():
    """Main page."""
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file upload - minimal version."""
    try:
        # Validate request size first
        valid_size, size_error = validate_request_size()
        if not valid_size:
            flash(size_error, 'error')
            return redirect(request.url)
            
        if 'file' not in request.files:
            flash('No file selected', 'error')
            return redirect(request.url)
        
        file = request.files['file']
        
        if file.filename == '':
            flash('No file selected', 'error')
            return redirect(request.url)
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            
            try:
                # Use very small temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                    file.save(temp_file.name)
                    temp_filepath = temp_file.name
                
                # Minimal analysis
                result = minimal_ppt_analysis(temp_filepath)
                
                # Clean up immediately
                os.unlink(temp_filepath)
                
                return jsonify({
                    'status': 'success',
                    'message': 'File analyzed successfully',
                    'result': result
                })
                
            except Exception as e:
                if 'temp_filepath' in locals() and os.path.exists(temp_filepath):
                    os.unlink(temp_filepath)
                return jsonify({'error': f'Analysis failed: {str(e)}'}), 500
            
        else:
            flash('Invalid file type. Please upload a .pptx or .ppt file under 500KB', 'error')
            return redirect(request.url)
    
    except Exception as e:
        return jsonify({'error': f'Upload failed: {str(e)}'}), 500

@app.route('/health')
def health_check():
    """Health check."""
    return jsonify({
        'status': 'healthy',
        'service': 'Minimal PPT Analyzer',
        'max_file_size_kb': MAX_FILE_SIZE / 1024,
        'version': '1.0.0-minimal'
    })

@app.route('/limits')
def get_limits():
    """Get current limits."""
    return jsonify({
        'max_file_size_bytes': MAX_FILE_SIZE,
        'max_file_size_kb': MAX_FILE_SIZE / 1024,
        'allowed_extensions': list(ALLOWED_EXTENSIONS),
        'note': 'Extremely conservative limits for Vercel serverless'
    })

# Error handlers
@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum: 500KB for Vercel serverless.'}), 413

@app.errorhandler(500)
def internal_error(e):
    return jsonify({'error': 'Internal server error', 'details': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)
